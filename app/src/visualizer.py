'''
Created on Sep 6, 2017

@author: Hossein

This file handles all the visualizations. It assumes that the status are ready and it
contains the function to generate the jpg's
'''
import os
from utils import get_test_path
import numpy as np
import matplotlib
# Force matplotlib to not use any Xwindows backend.
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pylab import gca
from datetime import datetime as dt
from pptx import Presentation
from pptx.util import Inches
import chart_plot


IMG_QUALITY = 400 # dpi

# Currently supported color themes: Grey, Purple, Blue, Green, Orange, Red
def generate_progress_chart(raised, goal, file_name= None, theme_color= 'Purple'):
    plt = chart_plot.prgress_chart(raised, goal, theme_color)
    plt.savefig(file_name, dpi=IMG_QUALITY)

def generate_bar_chart(label_value_pairs, file_name = None):
  '''
  generates a bar chart based on values, labels the x-axis with labels,
  and saves it to a file.
  input
    :list : a list of lables of the chart. These could be the
            divisions of organization
    :list : a list of values for each lable
    :str : a file_name to save the final result in
  '''
  label_value_pairs = sorted(label_value_pairs)
  n = len(label_value_pairs)
  x = np.arange(n)
  labels = list(zip(*label_value_pairs))[0]
  vals = list(zip(*label_value_pairs))[1]
  y = np.array(vals)
  # Clear and create the plot
  plt.clf()
  plt.bar(x, y, facecolor='#ff9999', edgecolor='white', ) # 9999ff ff9999
  # Write values
  for x1,y1 in zip(x, y):
      plt.text(x1, y1, '$%.0f' % y1, ha='center', va= 'bottom', fontsize = 30)
  # Adjust the y size so that there is room for writings
  plt.ylim(0, 1.2 * max(y) )
  plt.xticks(x, labels, fontsize = 22)
  # no y axis ticks
  plt.yticks([])
  # plt.title("Fundraising Showdown", fontsize = 30)
  # Remove all the spines except the bottom one
  # [trick here](https://stackoverflow.com/questions/18968024/how-to-remove-axis-in-pyplot-bar)
  for loc, spine in plt.axes().axes.spines.items():
    if loc != 'bottom':
      # no spine if it is not at the bottom
      spine.set_color('none')
  # Save the plot
  plt.savefig(file_name, dpi=IMG_QUALITY)


def generate_time_series(time_value_pairs, file_name):
  '''
  generates a xtime series chart based on values, labels the x-axis with xtime,
  and saves it to a file.
  input
    :list : a list of time_value pairs, e.g., [('YYYY-MM-DD', val_float), (...), ...]
    :str : a file_name to save the final result in
  '''
  times = list(zip(*time_value_pairs))[0]
  vals = list(zip(*time_value_pairs))[1]
  x = np.array([dt.strptime(t, "%Y-%m-%d") for t in times])
  y = np.array(vals)
  dates = matplotlib.dates.date2num(x)
  #print(dates)

  plt.clf()
  plt.plot_date(dates, y, fmt="bo-")
  for loc, spine in plt.axes().axes.spines.items():
    if loc != 'bottom':
      # no spine if it is not at the bottom
      spine.set_color('none')
  #remove y ticks
  plt.yticks([])
  # change the font size of x asix:
  ax = gca()
  ax.xaxis.set_tick_params(labelsize=14)
  # write the data point values on them
  nlabels = 3
  xy_list = list(zip(x, y))
  if len(xy_list) > nlabels:
    index_list = []
    for i in range(nlabels):
      index_list.append(i * (len(xy_list) // nlabels))
    index_list.append(len(xy_list) - 1)
    for i in index_list:
      x1, y1 = xy_list[i]
      plt.text(x1, y1, '$%.0f' % y1, ha='right', va= 'bottom', fontsize = 20)
  else:
    for x1, y1 in xy_list :
      plt.text(x1, y1, '$%.0f' % y1, ha='right', va= 'bottom', fontsize = 20)

  plt.gcf().autofmt_xdate()
  plt.ylim(0.5 * min(y), 1.1 * max(y) )
  # plt.title("Total fund raised", fontsize=20)
  plt.savefig(file_name, dpi=IMG_QUALITY)

def generate_ppt(result_file_name,
                 highest_donations = None,
                 image_total_fund_raised = None,
                 image_fundraising_progress = None,
                 image_fund_by_division = None,
                 logo=None,
                 file_template_presentation = None):
  '''
  Create the PowerPoint
  '''
  #logo coordinates:
  logo_left = Inches(7)
  logo_top = Inches(0.3)
  logo_height = Inches(2)

  #get the seed presentation
  if file_template_presentation:
    prs = Presentation(file_template_presentation)
  else:
    prs = Presentation()
  # Create the title page
  '''
  title_slide_layout = prs.slide_layouts[0]
  slide = prs.slides.add_slide(title_slide_layout)
  title = slide.shapes.title
  subtitle = slide.placeholders[1]
  title.text = "Join Team Ericsson Today!"
  subtitle.text = "LightTheNight.ca"
  '''

  # Add the bullet slide
  if highest_donations:
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'Thanks to Our\nTop Supporters'
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    last_hd = None
    for window_day in highest_donations:
      highest_donation = window_day[2]

      if highest_donation and highest_donation != last_hd:
        last_hd = highest_donation
        p = tf.add_paragraph()
        p.text =  window_day[0] + ":"
        p.level = 0
        p = tf.add_paragraph()
        p.text = '{} donated ${} to {}.'\
                .format(highest_donation['supporter_name'], highest_donation['amount_dollar'],
                        highest_donation['team_member'])
        p.level = 1
        if highest_donation['message']:
          p = tf.add_paragraph()
          p.text = '{} said: "{}" '.format(highest_donation['supporter_name'], highest_donation['message'])
          p.level = 2
      # add logo:
      if logo:
        pic = slide.shapes.add_picture(logo, logo_left, logo_top, height=logo_height)


  # add image for historical data
  '''
  if image_total_fund_raised:
    blank_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'Total Fund Raised\nby Ericsson Ottawa'
    left = Inches(2.5)
    top = Inches(2)
    height = Inches(5.5)
    pic = slide.shapes.add_picture(image_total_fund_raised, left, top, height=height)
    if logo:
      pic = slide.shapes.add_picture(logo, logo_left, logo_top, height=logo_height)
    '''

  # add slide for thermometer and bar chart on the same slide
  if image_fundraising_progress or image_fund_by_division:
    blank_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'Our Fundrasing Status'

    if logo:
      pic = slide.shapes.add_picture(logo, logo_left, logo_top, height=logo_height)

  # add image for bar chart
  if image_fund_by_division:
    # title_shape = shapes.title
    # title_shape.text = 'Fundraising Showdown'
    left = Inches(6)
    top = Inches(2)
    height = Inches(5)
    pic = slide.shapes.add_picture(image_fund_by_division, left, top, height=height)

  # add the progress chart
  if image_fundraising_progress:
    left = Inches(0.1)
    top = Inches(2)
    height = Inches(5)
    pic = slide.shapes.add_picture(image_fundraising_progress, left, top, height=height)

  # Save the powerpoint
  prs.save(result_file_name)


if __name__ == "__main__":
  # generate the bar chart for tema competition
  file_name_by_divisions = os.path.join(get_test_path(), 'divisions.png')
  l = [('Warehouse', 326), ('Tower', 1000)]
  generate_bar_chart(l,  file_name_by_divisions)

  # generate the total fund raised as a function of time
  file_name_time_series = os.path.join(get_test_path(), 'time_series.png')
  time_vals = [('2017-09-01', 100), ('2017-09-02', 200), ('2017-09-04', 310), ('2017-09-05', 400),
        ('2017-09-06', 450), ('2017-09-07', 600), ('2017-09-08', 600), ('2017-09-09', 600),
        ('2017-09-10', 600), ('2017-09-12', 600)]
  generate_time_series(time_vals, file_name_time_series)

  # generating ppt presentation
  print("generating ppt presentation...")
  from utils import get_raw_data_path
  supporters_data = {'2017-09-08': {'Johanna Nicoletta': [{'supporter_name': 'ericssoncommunity', 'message': '', 'amount_dollar': '229.41', 'time': 'Jul 19, 2017 12:00 AM'}, {'supporter_name': 'ericssoncommunity', 'message': '', 'amount_dollar': '27.38', 'time': 'May 12, 2017 12:00 AM'}, {'supporter_name': 'Anonymous Donor', 'message': '', 'amount_dollar': '250.00', 'time': 'May 1, 2017 8:45 AM'}], 'Ericsson Activities': [{'supporter_name': 'Foosball Champion', 'message': '', 'amount_dollar': '78.00', 'time': 'Aug 30, 2017 9:35 PM'}, {'supporter_name': 'Bottle Drive -- Aug 8', 'message': '', 'amount_dollar': '194.00', 'time': 'Aug 9, 2017 10:07 AM'}], 'Hossein Seyedmehdi': [], 'Alireza Mirzaee': [{'supporter_name': 'Hossein', 'message': '', 'amount_dollar': '25.00', 'time': 'Aug 9, 2017 10:18 AM'}]}, '2017-09-05': {'Johanna Nicoletta': [{'supporter_name': 'ericssoncommunity', 'message': '', 'amount_dollar': '27.38', 'time': 'May 12, 2017 12:00 AM'}, {'supporter_name': 'Anonymous Donor', 'message': '', 'amount_dollar': '250.00', 'time': 'May 1, 2017 8:45 AM'}], 'Ericsson Activities': [{'supporter_name': 'Foosball Champion', 'message': '', 'amount_dollar': '78.00', 'time': 'Aug 30, 2017 9:35 PM'}, {'supporter_name': 'Bottle Drive -- Aug 8', 'message': '', 'amount_dollar': '194.00', 'time': 'Aug 9, 2017 10:07 AM'}], 'Hossein Seyedmehdi': [], 'Alireza Mirzaee': [{'supporter_name': 'Hossein', 'message': '', 'amount_dollar': '25.00', 'time': 'Aug 9, 2017 10:18 AM'}]}, '2017-09-09': {'Johanna Nicoletta': [{'supporter_name': 'ericssoncommunity', 'message': '', 'amount_dollar': '229.41', 'time': 'Jul 19, 2017 12:00 AM'}, {'supporter_name': 'ericssoncommunity', 'message': '', 'amount_dollar': '27.38', 'time': 'May 12, 2017 12:00 AM'}, {'supporter_name': 'Anonymous Donor', 'message': '', 'amount_dollar': '250.00', 'time': 'May 1, 2017 8:45 AM'}], 'Ericsson Activities': [{'supporter_name': 'Foosball Champion', 'message': '', 'amount_dollar': '78.00', 'time': 'Aug 30, 2017 9:35 PM'}, {'supporter_name': 'Bottle Drive -- Aug 8', 'message': '', 'amount_dollar': '194.00', 'time': 'Aug 9, 2017 10:07 AM'}], 'Hossein Seyedmehdi': [], 'Alireza Mirzaee': [{'supporter_name': 'Hossein', 'message': '', 'amount_dollar': '25.00', 'time': 'Aug 9, 2017 10:18 AM'}]}, '2017-09-06': {'Johanna Nicoletta': [{'supporter_name': 'ericssoncommunity', 'message': '', 'amount_dollar': '27.38', 'time': 'May 12, 2017 12:00 AM'}, {'supporter_name': 'Anonymous Donor', 'message': '', 'amount_dollar': '250.00', 'time': 'May 1, 2017 8:45 AM'}], 'Ericsson Activities': [{'supporter_name': 'Foosball Champion', 'message': '', 'amount_dollar': '78.00', 'time': 'Aug 30, 2017 9:35 PM'}, {'supporter_name': 'Bottle Drive -- Aug 8', 'message': '', 'amount_dollar': '194.00', 'time': 'Aug 9, 2017 10:07 AM'}], 'Hossein Seyedmehdi': [], 'Alireza Mirzaee': [{'supporter_name': 'Hossein', 'message': '', 'amount_dollar': '25.00', 'time': 'Aug 9, 2017 10:18 AM'}]}, '2017-09-03': {'Johanna Nicoletta': [{'supporter_name': 'ericssoncommunity', 'message': '', 'amount_dollar': '27.38', 'time': 'May 12, 2017 12:00 AM'}, {'supporter_name': 'Anonymous Donor', 'message': '', 'amount_dollar': '250.00', 'time': 'May 1, 2017 8:45 AM'}], 'Ericsson Activities': [{'supporter_name': 'Foosball Champion', 'message': '', 'amount_dollar': '78.00', 'time': 'Aug 30, 2017 9:35 PM'}, {'supporter_name': 'Bottle Drive -- Aug 8', 'message': '', 'amount_dollar': '194.00', 'time': 'Aug 9, 2017 10:07 AM'}], 'Hossein Seyedmehdi': [], 'Alireza Mirzaee': [{'supporter_name': 'Hossein', 'message': '', 'amount_dollar': '25.00', 'time': 'Aug 9, 2017 10:18 AM'}]}, '2017-09-07': {'Johanna Nicoletta': [{'supporter_name': 'ericssoncommunity', 'message': '', 'amount_dollar': '27.38', 'time': 'May 12, 2017 12:00 AM'}, {'supporter_name': 'Anonymous Donor', 'message': '', 'amount_dollar': '250.00', 'time': 'May 1, 2017 8:45 AM'}], 'Ericsson Activities': [{'supporter_name': 'Foosball Champion', 'message': '', 'amount_dollar': '78.00', 'time': 'Aug 30, 2017 9:35 PM'}, {'supporter_name': 'Bottle Drive -- Aug 8', 'message': '', 'amount_dollar': '194.00', 'time': 'Aug 9, 2017 10:07 AM'}], 'Hossein Seyedmehdi': [], 'Alireza Mirzaee': [{'supporter_name': 'Hossein', 'message': '', 'amount_dollar': '25.00', 'time': 'Aug 9, 2017 10:18 AM'}]}}

  file_name_pptx = os.path.join(get_test_path(), 'resutls.pptx')
  file_name_template_pptx = os.path.join(get_raw_data_path(), 'template.pptx')
  file_name_logo = os.path.join(get_raw_data_path(), 'logo_LTN.png')
  file_name_fundraising_progress = os.path.join(get_raw_data_path(), 'fundrasing_progress.jpg')
  highest_donations = [['Yesterday', 2, {'team_member': 'Hossein Seyedmehdi', 'amount_dollar': '20.00', 'message': 'Good Job!', 'time': 'May 2, 2017 8:45 AM', 'supporter_name': 'SomeONE'}], ['Last week', 7, None], ['Last 30 days', 30, {'team_member': 'Johanna Nicoletta', 'amount_dollar': '250.00', 'message': '', 'time': 'May 1, 2017 8:45 AM', 'supporter_name': 'Anonymous Donor'}], ['Overall in 2017', 365, {'team_member': 'Johanna Nicoletta', 'amount_dollar': '250.00', 'message': '', 'time': 'May 1, 2017 8:45 AM', 'supporter_name': 'Anonymous Donor'}]]
  generate_ppt(file_name_pptx,
               highest_donations= highest_donations,
               image_total_fund_raised = file_name_time_series,
               image_fund_by_division = file_name_by_divisions,
               image_fundraising_progress= file_name_fundraising_progress,
               logo = file_name_logo,
               file_template_presentation= file_name_template_pptx)

  print('Done!')
